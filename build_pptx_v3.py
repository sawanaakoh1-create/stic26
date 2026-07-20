# -*- coding: utf-8 -*-
"""
Génère AfriVoice_AI_Pitch_STIC26_v3.pptx

Version FINALE alignée sur la fiche de candidature :
- La fiche ne fixe PAS de langue pilote → slide 5 réoriente le récit vers
  "choix expert de la langue pilote" (avec l'expert linguistique à recruter).
- Le prototype en ligne est présenté comme un APERÇU DE VISION, pas un
  produit final — pour préparer les questions du jury en toute honnêteté.
- Slide 10 renforce les 3 recrutements listés dans la fiche.
- Chiffres financiers alignés à la fiche (An1: -20k, An2: +25k, An3: +130k)
  et utilisateurs (3k → 15k → 50k).

Présentatrice : AKOH N'DJARMA M. Sawanatou (Togo), finaliste sélectionnée.
Concours porté par le Burkina Faso. Présentation en ligne, 10 min strictes.
11 slides, ~55 s/slide en moyenne.

Usage :
    py -3 build_pptx_v3.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================================
# CHARTE GRAPHIQUE
# ============================================================================
OCRE       = RGBColor(0xC9, 0x7B, 0x2E)   # ocre Sahel
TERRACOTTA = RGBColor(0xB0, 0x45, 0x2A)   # terre cuite
VERT_SAHEL = RGBColor(0x2E, 0x6E, 0x4E)   # vert profond
NUIT       = RGBColor(0x1A, 0x1F, 0x2E)   # bleu nuit
SABLE      = RGBColor(0xF5, 0xEE, 0xE0)   # sable clair
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
GRIS       = RGBColor(0x55, 0x5A, 0x66)
JAUNE      = RGBColor(0xE8, 0xB8, 0x3D)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============================================================================
# HELPERS
# ============================================================================

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NUIT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=NUIT,
                bullet="•", font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def add_header_band(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), OCRE)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.10), VERT_SAHEL)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.65),
             title, size=28, bold=True, color=BLANC)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
                 subtitle, size=14, color=SABLE)

def add_footer(slide, page_num, total=11):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), NUIT)
    add_text(slide, Inches(0.4), Inches(7.19), Inches(8), Inches(0.3),
             "AfriVoice AI  ·  STIC'26  ·  Finale en ligne  ·  Vision & prototype",
             size=10, color=SABLE)
    add_text(slide, Inches(11.5), Inches(7.19), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=10, color=SABLE, align=PP_ALIGN.RIGHT)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================================
# BUILD
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ---------- SLIDE 1 : COUVERTURE ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NUIT)
add_rect(s, 0, 0, Inches(0.35), SLIDE_H, OCRE)
add_rect(s, 0, Inches(2.2), SLIDE_W, Inches(0.06), OCRE)
add_rect(s, 0, Inches(5.2), SLIDE_W, Inches(0.06), VERT_SAHEL)

add_text(s, Inches(0.9), Inches(0.55), Inches(11), Inches(0.5),
         "STIC'26 · SAHEL TECH INNOVATION CHALLENGE 2026",
         size=14, bold=True, color=JAUNE)
add_text(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.5),
         "Finale en ligne  ·  Présentation 10 minutes  ·  Vision + prototype",
         size=13, color=SABLE)

add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4),
         "AfriVoice AI", size=72, bold=True, color=BLANC)
add_text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.7),
         "L'IA vocale pour l'Afrique de l'Ouest — pour ceux qui ne savent ni lire, ni écrire.",
         size=22, color=OCRE)

add_text(s, Inches(0.9), Inches(5.5), Inches(11), Inches(0.4),
         "AKOH N'DJARMA M. Sawanatou   ·   Fondatrice   ·   Lomé, Togo",
         size=16, bold=True, color=BLANC)
add_text(s, Inches(0.9), Inches(5.95), Inches(11), Inches(0.4),
         "Finaliste sélectionnée — STIC'26 (Ouagadougou, Burkina Faso)",
         size=13, color=SABLE)
add_text(s, Inches(0.9), Inches(6.4), Inches(11), Inches(0.4),
         "Prototype de vision : https://stic26.vercel.app",
         size=13, color=JAUNE)

# ---------- SLIDE 2 : LE PROBLÈME ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, SABLE)
add_header_band(s, "Le problème", "La barrière invisible du numérique en Afrique de l'Ouest")

def stat_card(slide, x, y, w, h, big, small):
    add_rect(slide, x, y, w, h, BLANC)
    add_rect(slide, x, y, w, Inches(0.1), OCRE)
    add_text(slide, x, y+Inches(0.55), w, Inches(1.2), big,
             size=44, bold=True, color=TERRACOTTA, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.2), y+Inches(1.9), w-Inches(0.4), Inches(1.2), small,
             size=14, color=NUIT, align=PP_ALIGN.CENTER)

stat_card(s, Inches(0.6), Inches(1.7), Inches(4),  Inches(3.2),
          "60 %+", "de la population\nutilise quotidiennement\nune langue locale")
stat_card(s, Inches(4.85), Inches(1.7), Inches(4), Inches(3.2),
          "65 %", "au Burkina Faso\nne lisent pas couramment\nle français (UNESCO)")
stat_card(s, Inches(9.1), Inches(1.7), Inches(3.7), Inches(3.2),
          "0", "assistant vocal IA\nlocal industriel disponible\naujourd'hui")

add_rect(s, Inches(0.6), Inches(5.2), Inches(12.2), Inches(1.5), NUIT)
add_text(s, Inches(0.9), Inches(5.4), Inches(11.6), Inches(1.2),
         "Elles ont un téléphone. Elles ont besoin de santé, d'agriculture,\n"
         "de finance, d'éducation. Mais le numérique leur parle en français,\n"
         "à l'écrit. Pour elles, c'est un mur.",
         size=17, color=BLANC, align=PP_ALIGN.CENTER)
add_footer(s, 2)

# ---------- SLIDE 3 : AWA ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BLANC)
add_rect(s, 0, 0, Inches(4.5), SLIDE_H, OCRE)
add_text(s, Inches(0.4), Inches(1.2), Inches(3.8), Inches(1.4),
         "Awa", size=90, bold=True, color=BLANC)
add_text(s, Inches(0.4), Inches(2.9), Inches(3.8), Inches(0.6),
         "42 ans", size=22, bold=True, color=SABLE)
add_text(s, Inches(0.4), Inches(3.4), Inches(3.8), Inches(0.6),
         "Commune rurale, Afrique de l'Ouest", size=16, color=SABLE)
add_text(s, Inches(0.4), Inches(3.9), Inches(3.8), Inches(0.6),
         "Mère de 4 enfants", size=18, color=SABLE)
add_text(s, Inches(0.4), Inches(4.7), Inches(3.8), Inches(0.6),
         "Ne sait pas lire.", size=20, bold=True, color=NUIT)

add_text(s, Inches(5.1), Inches(1.5), Inches(7.8), Inches(1),
         "Une histoire, un million de fois.", size=24, bold=True, color=TERRACOTTA)

add_bullets(s, Inches(5.1), Inches(2.5), Inches(7.8), Inches(2.5), [
    "Son fils a de la fièvre.",
    "Elle prend son téléphone. Elle abandonne devant les icônes.",
    "Elle n'a personne à qui demander.",
], size=17, color=NUIT)

add_rect(s, Inches(5.1), Inches(5.0), Inches(7.8), Inches(1.7), VERT_SAHEL)
add_text(s, Inches(5.4), Inches(5.15), Inches(7.2), Inches(1.5),
         "Avec AfriVoice AI\n"
         "elle appuie sur UN bouton, parle dans SA langue,\n"
         "et l'IA lui répond à voix haute.",
         size=17, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_footer(s, 3)

# ---------- SLIDE 4 : LA SOLUTION ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, SABLE)
add_header_band(s, "La solution",
                "Une infrastructure IA voice-first — 100 % pensée pour l'oral")

def pipe_block(slide, x, y, w, h, tag, title, desc, color):
    add_rect(slide, x, y, w, h, BLANC)
    add_rect(slide, x, y, w, Inches(0.55), color)
    add_text(slide, x, y+Inches(0.05), w, Inches(0.45), tag,
             size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.2), y+Inches(0.75), w-Inches(0.4), Inches(0.6), title,
             size=20, bold=True, color=NUIT, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.2), y+Inches(1.5), w-Inches(0.4), Inches(2.2), desc,
             size=13, color=GRIS, align=PP_ALIGN.CENTER)

y0 = Inches(1.85); h0 = Inches(3.8); w0 = Inches(3.7)
pipe_block(s, Inches(0.5), y0, w0, h0, "1 · ENTRÉE",  "ASR",
           "Reconnaissance vocale\n\nLa parole locale\nest transcrite en texte", OCRE)
pipe_block(s, Inches(4.8), y0, w0, h0, "2 · CERVEAU", "NLP",
           "Traitement du langage\n\nComprend l'intention\net formule la réponse", TERRACOTTA)
pipe_block(s, Inches(9.1), y0, w0, h0, "3 · SORTIE",  "TTS",
           "Synthèse vocale\n\nLa réponse est prononcée\navec une voix naturelle", VERT_SAHEL)

for cx in (Inches(4.5), Inches(8.8)):
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, cx, Inches(3.5),
                             Inches(0.35), Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb = NUIT
    arr.line.fill.background()

add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), NUIT)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.7),
         "Latence cible < 3 secondes bout en bout  ·  Fonctionne sur téléphones d'entrée de gamme",
         size=15, bold=True, color=JAUNE, align=PP_ALIGN.CENTER)
add_footer(s, 4)

# ---------- SLIDE 5 : CHOIX DE LA LANGUE PILOTE ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BLANC)
add_header_band(s, "Choix de la langue pilote",
                "Une décision experte — pas une décision marketing")

# Colonne gauche : critères
add_text(s, Inches(0.5), Inches(1.6), Inches(7.2), Inches(0.5),
         "4 critères objectifs de sélection", size=18, bold=True, color=TERRACOTTA)
add_bullets(s, Inches(0.5), Inches(2.15), Inches(7.2), Inches(3.5), [
    "Disponibilité et qualité des corpus vocaux existants.",
    "Nombre de locuteurs actifs dans la sous-région.",
    "Ancrage terrain de nos partenaires ONG / coopératives.",
    "Stratégie régionale et opportunités B2B / B2G.",
], size=15)

# Candidates
add_text(s, Inches(0.5), Inches(4.8), Inches(7.2), Inches(0.5),
         "Candidates naturelles", size=16, bold=True, color=VERT_SAHEL)
add_text(s, Inches(0.5), Inches(5.35), Inches(7.2), Inches(1.5),
         "Mooré · Dioula · Fulfuldé · Wolof · Bambara\n\n"
         "Le prototype en ligne illustre le concept avec\n"
         "des échantillons en Mooré (respect du pays hôte).",
         size=14, color=NUIT)

# Encadré droite : arbitrage expert
add_rect(s, Inches(8.2), Inches(1.6), Inches(4.7), Inches(4.9), OCRE)
add_text(s, Inches(8.4), Inches(1.8), Inches(4.3), Inches(0.5),
         "Décision arbitrée par l'expert", size=16, bold=True, color=BLANC)
add_text(s, Inches(8.4), Inches(2.4), Inches(4.3), Inches(4.0),
         "La fiche de candidature ne fixe pas la\n"
         "langue de démarrage : c'est un choix\n"
         "méthodologique assumé.\n\n"
         "Le choix définitif sera arrêté avec\n"
         "l'Expert Linguistique & Data Manager\n"
         "que je recruterai dès l'obtention du\n"
         "financement (profil listé dans la fiche,\n"
         "section « Équipe »).\n\n"
         "Fixer la langue avant l'expert, ce serait\n"
         "mettre la charrue avant les bœufs.",
         size=12, color=BLANC)
add_footer(s, 5)

# ---------- SLIDE 6 : DÉMO LIVE — APERÇU DE VISION ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NUIT)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), OCRE)
add_text(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.65),
         "Démo live — aperçu de la vision", size=28, bold=True, color=BLANC)
add_text(s, Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
         "Prototype construit en solo · illustre l'expérience utilisateur cible",
         size=14, color=SABLE)
add_rect(s, 0, Inches(1.15), SLIDE_W, Inches(0.10), VERT_SAHEL)

# Bandeau honnêteté
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.6), TERRACOTTA)
add_text(s, Inches(0.7), Inches(1.62), Inches(12), Inches(0.5),
         "APERÇU DE VISION — pas encore le modèle IA de production. "
         "Le prototype prouve le parcours utilisateur.",
         size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

# QR placeholder
add_rect(s, Inches(9.0), Inches(2.4), Inches(3.8), Inches(3.4), BLANC)
add_text(s, Inches(9.0), Inches(3.5), Inches(3.8), Inches(0.6),
         "[ QR CODE ]", size=22, bold=True, color=NUIT, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.0), Inches(4.1), Inches(3.8), Inches(0.5),
         "stic26.vercel.app", size=14, color=GRIS, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.0), Inches(5.9), Inches(3.8), Inches(0.4),
         "Scannez ou cliquez sur le lien du chat",
         size=12, color=SABLE, align=PP_ALIGN.CENTER)

# 4 scénarios à gauche
add_text(s, Inches(0.5), Inches(2.4), Inches(8), Inches(0.5),
         "4 scénarios · 1 seule interface : la voix",
         size=18, bold=True, color=OCRE)

def scen(slide, x, y, ico, title, sub):
    add_rect(slide, x, y, Inches(4), Inches(1.2), BLANC)
    add_rect(slide, x, y, Inches(0.5), Inches(1.2), TERRACOTTA)
    add_text(slide, x+Inches(0.7), y+Inches(0.15), Inches(3.2), Inches(0.5),
             f"{ico}  {title}", size=14, bold=True, color=NUIT)
    add_text(slide, x+Inches(0.7), y+Inches(0.65), Inches(3.2), Inches(0.5),
             sub, size=11, color=GRIS)

scen(s, Inches(0.5), Inches(3.0), "🩺", "Santé maternelle", "Awa · commune rurale")
scen(s, Inches(4.6), Inches(3.0), "🌾", "Agriculture",      "Rasmane · cultivateur")
scen(s, Inches(0.5), Inches(4.3), "💰", "Inclusion finance","Salif · entrepreneur")
scen(s, Inches(4.6), Inches(4.3), "📚", "Éducation",        "Fatimata · commerçante")

add_rect(s, Inches(0.5), Inches(5.7), Inches(8), Inches(1.15), VERT_SAHEL)
add_text(s, Inches(0.7), Inches(5.8), Inches(7.6), Inches(1.0),
         "Je lance la démo Santé en direct.\n"
         "Aucun clavier. Aucun écran à lire. Juste la voix.",
         size=14, bold=True, color=BLANC)
add_footer(s, 6)

# ---------- SLIDE 7 : INNOVATION ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, SABLE)
add_header_band(s, "Notre innovation",
                "Ce que ni Google, ni OpenAI, ni Meta ne font sur les langues à faible ressource")

def inno_card(slide, x, y, w, h, num, title, desc, color):
    add_rect(slide, x, y, w, h, BLANC)
    add_rect(slide, x, y, Inches(0.7), h, color)
    add_text(slide, x, y+Inches(0.3), Inches(0.7), Inches(1), num,
             size=40, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.9), y+Inches(0.35), w-Inches(1.1), Inches(0.55),
             title, size=18, bold=True, color=NUIT)
    add_text(slide, x+Inches(0.9), y+Inches(0.95), w-Inches(1.1), Inches(2.5),
             desc, size=13, color=GRIS)

y0 = Inches(1.85); h0 = Inches(3.8); w0 = Inches(3.9)
inno_card(s, Inches(0.5), y0, w0, h0, "1", "Corpus terrain",
          "À collecter avec des locuteurs\nnatifs rémunérés — encadré par\nl'expert linguistique à recruter.", OCRE)
inno_card(s, Inches(4.7), y0, w0, h0, "2", "ASR + TTS fine-tunés",
          "Modèles adaptés aux langues\nà faible ressource — portés par\nle CTO Big Data à recruter.", TERRACOTTA)
inno_card(s, Inches(8.9), y0, w0, h0, "3", "UX sans écran",
          "DÉJÀ démontré aujourd'hui\nvia le prototype en ligne :\nappuyer, parler, écouter.", VERT_SAHEL)

add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), NUIT)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.7),
         "Objectif : 1er assistant vocal IA end-to-end contextualisé pour l'Afrique de l'Ouest.",
         size=15, bold=True, color=JAUNE, align=PP_ALIGN.CENTER)
add_footer(s, 7)

# ---------- SLIDE 8 : MARCHÉ + BUSINESS MODEL ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BLANC)
add_header_band(s, "Marché & modèle économique",
                "3 segments · 3 flux de revenus · 1 seul moteur technologique")

def seg(slide, x, y, w, h, tag, who, revenue, color):
    add_rect(slide, x, y, w, h, BLANC, line=color)
    add_rect(slide, x, y, w, Inches(0.55), color)
    add_text(slide, x, y+Inches(0.05), w, Inches(0.45), tag,
             size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.2), y+Inches(0.75), w-Inches(0.4), Inches(1.5),
             who, size=13, color=NUIT)
    add_rect(slide, x+Inches(0.2), y+Inches(2.6), w-Inches(0.4), Inches(0.03), color)
    add_text(slide, x+Inches(0.2), y+Inches(2.75), w-Inches(0.4), Inches(1.5),
             revenue, size=13, bold=True, color=color)

y0 = Inches(1.85); h0 = Inches(4.2); w0 = Inches(3.9)
seg(s, Inches(0.5), y0, w0, h0, "B2C",
    "Populations non alphabétisées,\nagriculteurs,\nmicro-entrepreneurs.",
    "Freemium\nAccès de base gratuit\n+ options premium", OCRE)
seg(s, Inches(4.7), y0, w0, h0, "B2B",
    "Banques & microfinances,\ncliniques (e-santé),\nplateformes e-commerce.",
    "Licences API\nFacturation à l'appel\nou au forfait", TERRACOTTA)
seg(s, Inches(8.9), y0, w0, h0, "B2G / ONG",
    "Administrations, agences\nde santé, ONG humanitaires,\ncoopératives.",
    "Contrats-cadres\nCampagnes vocales\nde masse", VERT_SAHEL)
add_footer(s, 8)

# ---------- SLIDE 9 : TRACTION + ROADMAP ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, SABLE)
add_header_band(s, "Traction & roadmap",
                "Ce qui existe aujourd'hui — ce que le financement permet de construire")

# Traction à gauche
add_text(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.5),
         "Aujourd'hui", size=20, bold=True, color=TERRACOTTA)

def kpi(slide, x, y, big, small):
    add_rect(slide, x, y, Inches(2.9), Inches(1.4), BLANC)
    add_text(slide, x, y+Inches(0.05), Inches(2.9), Inches(0.7), big,
             size=30, bold=True, color=OCRE, align=PP_ALIGN.CENTER)
    add_text(slide, x, y+Inches(0.85), Inches(2.9), Inches(0.5), small,
             size=10, color=GRIS, align=PP_ALIGN.CENTER)

kpi(s, Inches(0.5), Inches(2.2), "Vision",  "documentée, présentée\ndans la fiche STIC'26")
kpi(s, Inches(3.5), Inches(2.2), "Prototype","en ligne, fonctionnel\n(stic26.vercel.app)")
kpi(s, Inches(0.5), Inches(3.75), "< 3 s",  "latence cible visée\nsur réseau 3G/4G")
kpi(s, Inches(3.5), Inches(3.75), "95 %",   "précision ASR cible\n(WER — fiche)")

# Roadmap à droite
add_text(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
         "Roadmap 3 ans (fiche)", size=20, bold=True, color=VERT_SAHEL)

def road(slide, x, y, year, txt, color):
    add_rect(slide, x, y, Inches(1.2), Inches(1.0), color)
    add_text(slide, x, y+Inches(0.2), Inches(1.2), Inches(0.7), year,
             size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_rect(slide, x+Inches(1.2), y, Inches(4.7), Inches(1.0), BLANC)
    add_text(slide, x+Inches(1.4), y+Inches(0.1), Inches(4.5), Inches(0.85),
             txt, size=12, color=NUIT)

road(s, Inches(7.0), Inches(2.2), "An 1",
     "Recrutement équipe · choix langue pilote\n"
     "Corpus + MVP industrialisé · 3 000 utilisateurs", OCRE)
road(s, Inches(7.0), Inches(3.35), "An 2",
     "2e langue locale · premiers contrats B2B / B2G\n"
     "15 000 utilisateurs actifs", TERRACOTTA)
road(s, Inches(7.0), Inches(4.5), "An 3",
     "3e et 4e langues · API publique · expansion sous-région\n"
     "50 000 utilisateurs actifs", VERT_SAHEL)

# Ligne financière (chiffres alignés à la fiche)
add_rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.0), NUIT)
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
         "Projections financières (USD — fiche)",
         size=13, bold=True, color=JAUNE)
add_text(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.5),
         "An 1 : −20 000  (amorçage)     ·     An 2 : +25 000  (point mort)     ·     An 3 : +130 000  (rentable)",
         size=14, bold=True, color=BLANC)
add_footer(s, 9)

# ---------- SLIDE 10 : ÉQUIPE + ASK (LA SLIDE CLÉ) ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BLANC)
add_header_band(s, "Équipe & demande au jury",
                "La vision est portée — l'équipe reste à constituer")

# Bloc gauche : équipe
add_text(s, Inches(0.5), Inches(1.55), Inches(6.2), Inches(0.5),
         "L'équipe cible (fiche §5)", size=18, bold=True, color=TERRACOTTA)

add_rect(s, Inches(0.5), Inches(2.1), Inches(6.2), Inches(0.9), OCRE)
add_text(s, Inches(0.7), Inches(2.2), Inches(5.9), Inches(0.8),
         "Aujourd'hui : moi seule\n"
         "Fondatrice · Lead Dev IA · Chercheuse terrain",
         size=12, bold=True, color=BLANC)

add_text(s, Inches(0.5), Inches(3.15), Inches(6.2), Inches(0.4),
         "3 profils experts à recruter dès financement :",
         size=13, bold=True, color=NUIT)

def profile(slide, x, y, num, title, desc, color):
    add_rect(slide, x, y, Inches(0.55), Inches(0.85), color)
    add_text(slide, x, y+Inches(0.15), Inches(0.55), Inches(0.55), num,
             size=20, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.7), y+Inches(0.05), Inches(5.2), Inches(0.4),
             title, size=13, bold=True, color=NUIT)
    add_text(slide, x+Inches(0.7), y+Inches(0.42), Inches(5.2), Inches(0.5),
             desc, size=11, color=GRIS)

profile(s, Inches(0.5), Inches(3.65), "1", "CTO — DevOps & Big Data",
        "Industrialisation des modèles + sécurité des données", TERRACOTTA)
profile(s, Inches(0.5), Inches(4.6), "2", "Expert Linguistique & Data Manager",
        "Corpus vocaux + choix de la langue pilote", VERT_SAHEL)
profile(s, Inches(0.5), Inches(5.55), "3", "Responsable Produit & Commercial",
        "Acquisition B2B / B2G + UX vocale", OCRE)

# Bloc droite : ask
add_text(s, Inches(7.0), Inches(1.55), Inches(6), Inches(0.5),
         "Notre demande au jury", size=18, bold=True, color=VERT_SAHEL)

def ask_line(slide, x, y, n, txt):
    add_rect(slide, x, y, Inches(0.5), Inches(0.5), VERT_SAHEL)
    add_text(slide, x, y+Inches(0.05), Inches(0.5), Inches(0.4), n,
             size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.65), y+Inches(0.08), Inches(5.4), Inches(0.5),
             txt, size=13, color=NUIT)

ask_line(s, Inches(7.0), Inches(2.15), "1", "Financement amorçage (infra + fine-tuning)")
ask_line(s, Inches(7.0), Inches(2.85), "2", "Budget de recrutement des 3 profils experts")
ask_line(s, Inches(7.0), Inches(3.55), "3", "Appui à la collecte de corpus vocaux terrain")
ask_line(s, Inches(7.0), Inches(4.25), "4", "Partenariats ONG, coopératives, mairies, télécoms")

add_rect(s, Inches(7.0), Inches(5.05), Inches(6), Inches(1.6), NUIT)
add_text(s, Inches(7.2), Inches(5.2), Inches(5.6), Inches(1.4),
         "Votre soutien ne finance pas juste une startup.\n"
         "Il finance la CONSTITUTION d'une équipe experte\n"
         "capable de mettre des millions de personnes\n"
         "non alphabétisées dans l'ère numérique.",
         size=12, bold=True, color=JAUNE, align=PP_ALIGN.CENTER)
add_footer(s, 10)

# ---------- SLIDE 11 : VISION + MERCI ----------
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NUIT)
add_rect(s, 0, 0, Inches(0.35), SLIDE_H, OCRE)

add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.6),
         "Notre vision", size=18, bold=True, color=JAUNE)

add_text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.5),
         "Un continent où plus jamais un paysan, une mère, un artisan\n"
         "ne sera exclu du numérique parce qu'il ne sait pas lire.",
         size=26, bold=True, color=BLANC)

add_text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.5),
         "AfriVoice AI, ce n'est pas encore une application finie.\n"
         "C'est une vision documentée, un prototype honnête,\n"
         "et une fondatrice prête à s'entourer des meilleurs pour la réaliser.",
         size=17, color=OCRE)

# bande merci
add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(1.9), OCRE)
add_text(s, Inches(0.9), Inches(5.75), Inches(7), Inches(0.7),
         "Merci.", size=44, bold=True, color=BLANC)
add_text(s, Inches(0.9), Inches(6.55), Inches(8), Inches(0.4),
         "AKOH N'DJARMA M. Sawanatou · Lomé, Togo",
         size=14, bold=True, color=BLANC)
add_text(s, Inches(0.9), Inches(6.95), Inches(8), Inches(0.4),
         "Prototype de vision : https://stic26.vercel.app",
         size=13, color=SABLE)

# QR à droite dans la bande
add_rect(s, Inches(10.8), Inches(5.75), Inches(2.0), Inches(1.6), BLANC)
add_text(s, Inches(10.8), Inches(6.35), Inches(2.0), Inches(0.5),
         "[ QR ]", size=16, bold=True, color=NUIT, align=PP_ALIGN.CENTER)

# ============================================================================
# SAVE
# ============================================================================
out = Path(__file__).parent / "AfriVoice_AI_Pitch_STIC26_v3.pptx"
prs.save(out)
print(f"OK -> {out}")
print(f"Slides : {len(prs.slides)}")
